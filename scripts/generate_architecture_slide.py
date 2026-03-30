from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BG = RGBColor(6, 17, 28)
PANEL = RGBColor(10, 23, 38)
PANEL_ALT = RGBColor(12, 31, 49)
PANEL_SOFT = RGBColor(15, 38, 60)
TEXT = RGBColor(236, 245, 255)
MUTED = RGBColor(191, 209, 229)
CYAN = RGBColor(125, 230, 255)
BLUE = RGBColor(119, 164, 255)
GREEN = RGBColor(95, 214, 154)
AMBER = RGBColor(255, 208, 122)
ROSE = RGBColor(255, 165, 165)
LINE = RGBColor(67, 103, 145)


def set_run_style(run, *, size: int, color: RGBColor, bold: bool = False) -> None:
    font = run.font
    font.name = "Aptos"
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color


def add_textbox(slide, left, top, width, height, text, *, size=16, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run_style(run, size=size, color=color, bold=bold)
    return shape


def add_panel(slide, left, top, width, height, *, title, subtitle=None, fill=PANEL, title_color=TEXT, title_size=18):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.2)

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = title
    set_run_style(r1, size=title_size, color=title_color, bold=True)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(3)
        r2 = p2.add_run()
        r2.text = subtitle
        set_run_style(r2, size=10, color=MUTED)
    return shape


def add_bullets(slide, left, top, width, height, lines, *, size=11, color=MUTED):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.bullet = True
        p.level = 0
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line
        set_run_style(run, size=size, color=color)
    return box


def connect(slide, x1, y1, x2, y2, *, color=LINE, width=1.8, begin_arrow=False, end_arrow=True):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    if begin_arrow:
        line.line.begin_arrowhead = True
    if end_arrow:
        line.line.end_arrowhead = True
    return line


def add_flow_node(slide, left, top, width, height, *, label, detail, fill=PANEL_ALT, accent=CYAN):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.3)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(7)
    tf.margin_bottom = Pt(7)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = label
    set_run_style(r1, size=13, color=TEXT, bold=True)

    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = detail
    set_run_style(r2, size=9, color=MUTED)
    return shape


def build_slide() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    add_textbox(
        slide,
        Inches(0.35),
        Inches(0.18),
        Inches(7.2),
        Inches(0.45),
        "K8sWhisperer End-to-End Architecture",
        size=24,
        color=TEXT,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.38),
        Inches(0.60),
        Inches(7.8),
        Inches(0.28),
        "Real minikube incident signals -> LangGraph safety loop -> Slack HITL -> audit trail -> optional Soroban attestation",
        size=11,
        color=MUTED,
    )

    add_panel(
        slide,
        Inches(0.35),
        Inches(1.05),
        Inches(2.15),
        Inches(2.55),
        title="1. Signal Sources",
        subtitle="Live cluster inputs",
        fill=PANEL,
        title_color=CYAN,
    )
    add_bullets(
        slide,
        Inches(0.47),
        Inches(1.62),
        Inches(1.9),
        Inches(1.8),
        [
            "Namespace-scoped pods, events, deployments",
            "Optional read-only node snapshot",
            "Optional Prometheus metrics",
            "Demo workloads: CrashLoop, OOM, Pending",
        ],
    )

    add_panel(
        slide,
        Inches(0.35),
        Inches(4.05),
        Inches(2.15),
        Inches(2.10),
        title="Interfaces",
        subtitle="Operator and external surfaces",
        fill=PANEL,
        title_color=CYAN,
    )
    add_bullets(
        slide,
        Inches(0.47),
        Inches(4.62),
        Inches(1.9),
        Inches(1.3),
        [
            "FastAPI: health, incidents, audit, attestation",
            "Slack interactive approval callback",
            "React operator console for proofs",
            "Typed MCP servers: kubectl, Slack, Prometheus",
        ],
    )

    add_panel(
        slide,
        Inches(2.78),
        Inches(1.05),
        Inches(6.15),
        Inches(2.95),
        title="2. LangGraph Control Loop",
        subtitle="Shared typed state, persistent checkpoints, exact resume by incident_id",
        fill=PANEL_SOFT,
        title_color=BLUE,
    )

    node_w = Inches(0.78)
    node_h = Inches(0.92)
    start_x = Inches(3.0)
    top_y = Inches(1.92)
    step = Inches(0.82)
    labels = [
        ("observe", "pods/events"),
        ("detect", "hybrid LLM\n+ heuristic"),
        ("diagnose", "logs + evidence"),
        ("plan", "action + blast"),
        ("safety", "auto vs HITL"),
        ("execute", "restart/patch\nor guidance"),
        ("explain", "summary + log"),
    ]
    centers = []
    for idx, (label, detail) in enumerate(labels):
        left = start_x + step * idx
        add_flow_node(slide, left, top_y, node_w, node_h, label=label, detail=detail)
        centers.append((left + node_w / 2, top_y + node_h / 2))
        if idx:
            prev_x, prev_y = centers[idx - 1]
            cur_x, cur_y = centers[idx]
            connect(slide, prev_x + Inches(0.32), prev_y, cur_x - Inches(0.32), cur_y)

    add_textbox(slide, Inches(3.15), Inches(3.22), Inches(5.35), Inches(0.42), "Core incident paths: CrashLoopBackOff auto-remediation, OOMKilled HITL recommendation/patch path, PendingPod guidance-only", size=10, color=MUTED)

    add_panel(
        slide,
        Inches(9.18),
        Inches(1.05),
        Inches(3.78),
        Inches(2.95),
        title="3. Safety + Human Gate",
        subtitle="Safe automation, not blind automation",
        fill=PANEL,
        title_color=AMBER,
    )
    add_bullets(
        slide,
        Inches(9.34),
        Inches(1.62),
        Inches(3.4),
        Inches(1.9),
        [
            "Low-blast-radius plans can auto-approve",
            "Riskier plans pause with interrupt()",
            "Slack button click resumes the exact graph thread",
            "Default RBAC stays namespace-scoped",
            "Destructive actions stay denylisted by default",
        ],
        size=10,
    )

    connect(slide, Inches(8.95), Inches(2.45), Inches(9.18), Inches(2.45), color=AMBER)

    add_panel(
        slide,
        Inches(2.78),
        Inches(4.25),
        Inches(4.35),
        Inches(1.92),
        title="4. State, Query, and Audit",
        subtitle="Everything is explainable and inspectable",
        fill=PANEL,
        title_color=GREEN,
    )
    add_bullets(
        slide,
        Inches(2.94),
        Inches(4.82),
        Inches(4.0),
        Inches(1.05),
        [
            "Checkpoint store preserves pending approvals across restarts",
            "JSONL audit log stores diagnosis, decision, evidence, and result",
            "FastAPI summary endpoints merge runtime + audit for operators",
        ],
        size=10,
    )

    add_panel(
        slide,
        Inches(7.38),
        Inches(4.25),
        Inches(2.55),
        Inches(1.92),
        title="5. Frontend Desk",
        subtitle="Separate from control loop",
        fill=PANEL,
        title_color=CYAN,
    )
    add_bullets(
        slide,
        Inches(7.52),
        Inches(4.82),
        Inches(2.25),
        Inches(0.95),
        [
            "Browse incidents and audit timeline",
            "Anchor proof record",
            "Verify canonical hash",
        ],
        size=10,
    )

    add_panel(
        slide,
        Inches(10.18),
        Inches(4.25),
        Inches(2.78),
        Inches(1.92),
        title="6. Optional Soroban",
        subtitle="Bonus path only",
        fill=PANEL,
        title_color=ROSE,
    )
    add_bullets(
        slide,
        Inches(10.33),
        Inches(4.82),
        Inches(2.45),
        Inches(0.98),
        [
            "Backend anchors incident hash on Soroban",
            "Frontend verifies proof",
            "Not part of remediation execution",
        ],
        size=10,
    )

    connect(slide, Inches(6.20), Inches(4.05), Inches(5.05), Inches(4.25), color=GREEN, width=1.4)
    connect(slide, Inches(7.13), Inches(5.20), Inches(7.38), Inches(5.20), color=CYAN, width=1.4)
    connect(slide, Inches(9.93), Inches(5.20), Inches(10.18), Inches(5.20), color=ROSE, width=1.4)

    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(6.48), Inches(12.6), Inches(0.58))
    footer.fill.solid()
    footer.fill.fore_color.rgb = PANEL_ALT
    footer.line.color.rgb = LINE
    footer.line.width = Pt(1.0)
    tf = footer.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Key message: K8sWhisperer runs on a real minikube cluster, uses a persistent LangGraph workflow for safe automation, gates risky actions through Slack HITL, writes a full audit trail, and keeps Soroban attestation isolated as optional proof depth."
    set_run_style(run, size=11, color=TEXT, bold=False)

    return prs


if __name__ == "__main__":
    presentation = build_slide()
    output = "architecture_end_to_end.pptx"
    presentation.save(output)
    print(output)
